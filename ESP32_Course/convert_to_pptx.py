import os
from pptx import Presentation
from pptx.util import Inches
import docx

def get_mapping():
    return {
        "NTC_Thermistor_10k.jpg": "NTC Thermistor",
        "PT100_MAX31865.jpg": "PT100",
        "TMP36.jpg": "TMP36",
        "DS18B20.jpg": "DS18B20",
        "SHT31.jpg": "SHT31",
        "MAX6675_thermocouple.jpg": "MAX6675",
        "MLX90614_infrared_temperature_sensor.jpg": "MLX90614",
        "SHTC3_humidity_sensor.jpg": "SHTC3",
        "DHT22_humidity_sensor.jpg": "DHT22",
        "BME280_sensor_module.jpg": "BME280",
        "MQ-135_gas_sensor_module.jpg": "MQ-135",
        "MH-Z19B_CO2_sensor.jpg": "MH-Z19B",
        "SGP30_VOC_sensor.jpg": "SGP30",
        "BME680_gas_sensor.jpg": "BME680",
        "PMS5003_dust_sensor.jpg": "PMS5003",
        "GL5528_LDR_sensor.jpg": "GL5528",
        "BH1750_light_sensor.jpg": "BH1750",
        "TSL2561_light_sensor.jpg": "TSL2561",
        "TCS34725_color_sensor.jpg": "TCS3472",
        "APDS9999_sensor_module.jpg": "APDS9999",
        "VEML6075_UV_sensor.jpg": "VEML6075",
        "AS7262_spectral_sensor.jpg": "AS7262",
        "Capacitive_soil_moisture_sensor_v2.0.jpg": "Capacitive",
        "DFRobot_EC_sensor.jpg": "DFRobot",
        "Gravity_pH_sensor.jpg": "Gravity pH",
        "MPU6050_module.jpg": "MPU6050",
        "MPU9250_module.jpg": "MPU9250",
        "ADXL345_accelerometer.jpg": "ADXL345",
        "HC-SR04_ultrasonic_sensor.jpg": "HC-SR04",
        "VL53L0X_ToF_sensor.jpg": "VL53L0X",
        "GP2Y0A21_IR_distance_sensor.jpg": "GP2Y0A21",
        "BMP280_pressure_sensor.jpg": "BMP280",
        "HX711_load_cell_amplifier.jpg": "HX711",
        "FSR402_force_sensor.jpg": "FSR402",
        "A3144_hall_sensor.jpg": "A3144",
        "ACS712_current_sensor.jpg": "ACS712",
        "INA219_current_sensor.jpg": "INA219",
        "MAX17048_fuel_gauge.jpg": "MAX17048",
        "OV2640_camera_module.jpg": "OV2640",
        "FLIR_Lepton_thermal_camera.jpg": "FLIR",
        "Analog_microphone_module.jpg": "아날로그 마이크",
        "INMP441_MEMS_microphone.jpg": "INMP441",
        "TTP223_touch_sensor.jpg": "TTP223",
        "ESP32_capacitive_touch.jpg": "ESP32 터치핀",
        "YF-S201_water_flow_sensor.jpg": "YF-S201",
        "Float_switch_water_level_sensor.jpg": "Float switch",
        "Turbidity_sensor_module.jpg": "Turbidity"
    }

def main():
    prs = Presentation()
    
    # layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    blocks = []
    current_block = []
    
    with open('docx_dump.txt', 'r', encoding='utf-8') as f:
        lines = f.readlines()
        
    for line in lines:
        txt = line.strip('\n').strip()
        if not txt:
            if current_block:
                blocks.append(current_block)
                current_block = []
        else:
            current_block.append(txt)
            
    if current_block:
        blocks.append(current_block)
        
    mapping = get_mapping()
    current_section = "ESP32 Sensor Catalog"
    
    for idx, block in enumerate(blocks):
        block_text = "\n".join(block)
        
        # Check if it's likely a section heading like '1. 온도 센서'
        if len(block) == 1 and block[0][0].isdigit() and '. ' in block[0]:
            current_section = block[0]
            continue
            
        # First block should be Title slide
        if idx == 0:
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = block[0]
            subtitle.text = "\n".join(block[1:])
            continue
            
        # Content slide
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = current_section
        
        # We manually position text box to give space for images
        content_box = slide.placeholders[1]
        content_box.text = block_text
        
        # Make the text box smaller (on the left)
        content_box.width = Inches(4.5)
        
        # Find matches
        matched_images = []
        for img_file, kw in mapping.items():
            if kw.lower() in block_text.lower():
                img_path = os.path.join('sensor_images', img_file)
                if os.path.exists(img_path):
                    matched_images.append(img_path)
                    
        # Add images
        num_imgs = len(matched_images)
        for i, img_path in enumerate(matched_images):
            # Positioning magic (right side)
            left = Inches(5.0)
            # Center vertically based on number of images
            # Total available vertical height is around 5.5 inches
            img_height = Inches(5.0 / max(num_imgs, 1))
            if img_height > Inches(3):
                img_height = Inches(3)
                
            top = Inches(1.5 + i * (img_height.inches + 0.2))
            try:
                slide.shapes.add_picture(img_path, left, top, height=img_height)
            except Exception as e:
                print(f"Error adding {img_path}: {e}")

    prs.save('ESP32_Sensor_Catalog.pptx')
    print("Successfully created ESP32_Sensor_Catalog.pptx")

if __name__ == '__main__':
    main()
